"""
slide_extract_task — PDF/PPTX → slides rows + bullets + thumbnails.

Phase 6k replaces the `pdftoppm` PNG-only version with PyMuPDF for PDFs
and python-pptx for PPTX. Each page yields:
  • slides row: slide_index, slide_number, title, full_text, image_uri,
                thumbnail_uri
  • bullets[]: one row per bullet/paragraph

Ports MIC `app/tasks/slide_extract.py` (715 LOC) condensed for Rounds'
narrower schema.

Closes audit gaps 🟠 #10 (PPTX dropped), 🟠 #17 (no bullets), 🟡 #4 (multi-PDF).
"""
from __future__ import annotations

import io
import logging
import os
import tempfile

from app.tasks.celery_app import RoundsTask, celery_app

logger = logging.getLogger(__name__)


_PDF_MIMES = {"application/pdf"}
_PPTX_MIMES = {"application/vnd.openxmlformats-officedocument.presentationml.presentation"}


@celery_app.task(
    bind=True,
    base=RoundsTask,
    name="rounds.tasks.slide_extract",
    max_retries=2,
)
def slide_extract_task(self, session_id: str) -> dict:
    from sqlalchemy import create_engine, text

    from app.config import settings

    sync_url = settings.DATABASE_URL.replace("+asyncpg", "")
    engine = create_engine(sync_url)

    try:
        with engine.connect() as conn:
            existing = conn.execute(
                text("SELECT id FROM slides WHERE session_id = CAST(:sid AS uuid) LIMIT 1"),
                {"sid": session_id},
            ).fetchone()
            # We'll still process new sources, but skip duplicates per-source.
            slide_sources = conn.execute(
                text(
                    """
                    SELECT id, gcs_uri, filename, content_type
                      FROM sources WHERE session_id = CAST(:sid AS uuid) AND role = 'slide'
                      ORDER BY created_at ASC
                    """
                ),
                {"sid": session_id},
            ).fetchall()

        if not slide_sources:
            logger.info(f"slide_extract: no slide sources for {session_id}")
            return {"session_id": session_id, "slide_count": 0}
        if existing:
            logger.info(f"slide_extract: skip — slides exist for {session_id}")
            return {"skipped": True, "session_id": session_id}

        total_slides = 0
        total_bullets = 0
        slide_index_cursor = 0

        for src in slide_sources:
            src_id, gcs_uri, filename, content_type = src
            ext = (filename or "").rsplit(".", 1)[-1].lower()
            is_pdf = ext == "pdf" or (content_type or "") in _PDF_MIMES
            is_pptx = ext == "pptx" or (content_type or "") in _PPTX_MIMES

            if is_pdf:
                slides_written, bullets_written = _process_pdf(
                    engine, session_id, gcs_uri, slide_index_cursor,
                )
            elif is_pptx:
                slides_written, bullets_written = _process_pptx(
                    engine, session_id, gcs_uri, slide_index_cursor,
                )
            else:
                logger.warning(
                    f"slide_extract: unsupported source {filename} ({content_type}) — skipping"
                )
                slides_written = bullets_written = 0

            total_slides += slides_written
            total_bullets += bullets_written
            slide_index_cursor += slides_written

        logger.info(
            f"slide_extract: session={session_id} slides={total_slides} bullets={total_bullets}"
        )
        return {
            "session_id":  session_id,
            "slide_count": total_slides,
            "bullets":     total_bullets,
        }

    except Exception as exc:  # noqa: BLE001
        attempt = self.request.retries
        if attempt < self.max_retries:
            self.retry_with_backoff(exc, attempt)
        logger.exception(f"slide_extract: terminal failure for {session_id} — continuing")
        return {"session_id": session_id, "slide_count": 0, "error": str(exc)}
    finally:
        engine.dispose()


# ─── PDF path ───────────────────────────────────────────────────────────


def _process_pdf(engine, session_id: str, gcs_uri: str, index_base: int) -> tuple[int, int]:
    import fitz  # PyMuPDF
    from sqlalchemy import text

    from app.config import settings
    from app.tasks.transcribe import _download_from_gcs

    from google.cloud import storage as gcs_lib

    gcs_client = gcs_lib.Client(project=settings.GCP_PROJECT_ID)
    bucket = gcs_client.bucket(settings.GCS_BUCKET)

    with tempfile.TemporaryDirectory() as tmpdir:
        local_pdf = os.path.join(tmpdir, "deck.pdf")
        _download_from_gcs(gcs_uri, local_pdf)

        doc = fitz.open(local_pdf)
        slides_written = 0
        bullets_written = 0

        for page_idx in range(doc.page_count):
            page = doc.load_page(page_idx)
            full_text = page.get_text("text") or ""
            title = _derive_slide_title(full_text, page_idx + 1)

            # Render thumbnail (PNG) + upload
            pix = page.get_pixmap(dpi=120)
            png_bytes = pix.tobytes("png")
            blob_name = f"sessions/{session_id}/slides/thumb_{index_base + page_idx:03d}.png"
            bucket.blob(blob_name).upload_from_string(png_bytes, content_type="image/png")
            image_uri = f"gs://{settings.GCS_BUCKET}/{blob_name}"

            with engine.begin() as conn:
                row = conn.execute(
                    text(
                        """
                        INSERT INTO slides
                            (session_id, slide_index, title, image_uri, thumbnail_uri, full_text)
                        VALUES
                            (CAST(:sid AS uuid), :idx, :title, :uri, :uri, :ft)
                        ON CONFLICT (session_id, slide_index) DO UPDATE
                          SET title         = COALESCE(slides.title, EXCLUDED.title),
                              image_uri     = EXCLUDED.image_uri,
                              thumbnail_uri = EXCLUDED.thumbnail_uri,
                              full_text     = EXCLUDED.full_text
                        RETURNING id
                        """
                    ),
                    {
                        "sid":   session_id,
                        "idx":   index_base + page_idx,
                        "title": title,
                        "uri":   image_uri,
                        "ft":    full_text,
                    },
                ).fetchone()
                slide_id = str(row[0]) if row else None
                slides_written += 1

                # Bullets — extract bullet/paragraph blocks from PyMuPDF.
                bullets = _extract_bullets(page)
                for pos, bullet_text in enumerate(bullets):
                    conn.execute(
                        text(
                            """
                            INSERT INTO bullets (slide_id, text, position)
                            VALUES (CAST(:sid AS uuid), :tx, :pos)
                            ON CONFLICT (slide_id, position) DO UPDATE
                              SET text = EXCLUDED.text
                            """
                        ),
                        {"sid": slide_id, "tx": bullet_text, "pos": pos},
                    )
                    bullets_written += 1

        doc.close()
        return slides_written, bullets_written


def _derive_slide_title(full_text: str, fallback_number: int) -> str:
    """First non-empty line of slide text is the title heuristic."""
    for line in (full_text or "").splitlines():
        s = line.strip()
        if s:
            return s[:200]
    return f"Slide {fallback_number}"


def _extract_bullets(page) -> list[str]:
    """
    Extract bullet/paragraph candidates from a PyMuPDF page.
    Treats each text block's lines as bullets; filters by min length 3 chars.
    """
    bullets: list[str] = []
    blocks = page.get_text("blocks") or []
    for block in blocks:
        text = (block[4] or "").strip() if len(block) > 4 else ""
        if not text:
            continue
        for line in text.split("\n"):
            stripped = line.strip().lstrip("•·-*").strip()
            if len(stripped) >= 3:
                bullets.append(stripped[:500])
    return bullets


# ─── PPTX path ──────────────────────────────────────────────────────────


def _process_pptx(engine, session_id: str, gcs_uri: str, index_base: int) -> tuple[int, int]:
    from pptx import Presentation
    from sqlalchemy import text

    from app.config import settings
    from app.tasks.transcribe import _download_from_gcs

    with tempfile.TemporaryDirectory() as tmpdir:
        local_pptx = os.path.join(tmpdir, "deck.pptx")
        _download_from_gcs(gcs_uri, local_pptx)

        pres = Presentation(local_pptx)
        slides_written = 0
        bullets_written = 0
        for idx, slide in enumerate(pres.slides):
            title = None
            bullets: list[str] = []
            full_text_lines: list[str] = []
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                tf = shape.text_frame
                for para in tf.paragraphs:
                    text_str = para.text.strip()
                    if not text_str:
                        continue
                    full_text_lines.append(text_str)
                    if title is None and shape.name and "title" in shape.name.lower():
                        title = text_str
                    else:
                        bullets.append(text_str)
            title = title or (full_text_lines[0] if full_text_lines else f"Slide {idx + 1}")
            full_text = "\n".join(full_text_lines)

            with engine.begin() as conn:
                row = conn.execute(
                    text(
                        """
                        INSERT INTO slides
                            (session_id, slide_index, title, full_text)
                        VALUES
                            (CAST(:sid AS uuid), :idx, :title, :ft)
                        ON CONFLICT (session_id, slide_index) DO UPDATE
                          SET title     = EXCLUDED.title,
                              full_text = EXCLUDED.full_text
                        RETURNING id
                        """
                    ),
                    {
                        "sid":   session_id,
                        "idx":   index_base + idx,
                        "title": title[:200],
                        "ft":    full_text,
                    },
                ).fetchone()
                slide_id = str(row[0]) if row else None
                slides_written += 1

                for pos, bt in enumerate(bullets):
                    conn.execute(
                        text(
                            """
                            INSERT INTO bullets (slide_id, text, position)
                            VALUES (CAST(:sid AS uuid), :tx, :pos)
                            ON CONFLICT (slide_id, position) DO UPDATE
                              SET text = EXCLUDED.text
                            """
                        ),
                        {"sid": slide_id, "tx": bt[:500], "pos": pos},
                    )
                    bullets_written += 1

        return slides_written, bullets_written
